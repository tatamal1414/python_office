from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.pptx')




prs1 = Presentation()
title_slide_layout = prs1.slide_layouts[1]
slide = prs1.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "こんにちは！"
subtitle.text = "サブタイトルです！"

prs1.save('test2.pptx')